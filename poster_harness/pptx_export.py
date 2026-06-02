from __future__ import annotations

import tempfile
from pathlib import Path
from typing import Any, Mapping

from PIL import Image

from .replace import paste_fit, replace_placeholders


def export_image_as_pptx(
    image_path: str | Path,
    out_path: str | Path,
    *,
    long_edge_inches: float = 15.0,
) -> Path:
    """Create a one-slide PPTX containing the poster image at full-slide size.

    This is the legacy flat container export: the whole poster is one picture.
    Prefer :func:`export_layered_pptx` when a placement spec and assets are
    available, because it keeps source figures as separate PowerPoint objects.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:  # pragma: no cover - dependency declared in pyproject
        raise RuntimeError("python-pptx is required to export PPTX posters") from exc

    src = Path(image_path)
    out = Path(out_path)
    width_px, height_px = _image_size(src)
    slide_w_in, slide_h_in = _slide_inches(width_px, height_px, long_edge_inches=long_edge_inches)

    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(
        str(src),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


def export_layered_pptx(
    *,
    base_image: str | Path,
    spec: Mapping[str, Any],
    asset_dir: str | Path,
    out_path: str | Path,
    scale: float = 1.0,
    long_edge_inches: float = 15.0,
) -> Path:
    """Create a one-slide PPTX with layout and real figures as separate layers.

    The generated poster layout is still a raster background because the image
    model did not produce editable PowerPoint vector/text objects.  However, the
    deterministic source figures are inserted as individual PowerPoint picture
    objects at their placement boxes, so they can be selected, moved, resized, or
    replaced independently of the layout background.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:  # pragma: no cover - dependency declared in pyproject
        raise RuntimeError("python-pptx is required to export PPTX posters") from exc

    base = Path(base_image)
    out = Path(out_path)
    width_px, height_px = _image_size(base)
    slide_w_in, slide_h_in = _slide_inches(width_px, height_px, long_edge_inches=long_edge_inches)

    out.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="poster_harness_pptx_") as tmp_dir:
        tmp = Path(tmp_dir)
        background = tmp / "layout-background.png"
        # Same deterministic cleanup/frame pass as PNG export, but without
        # pasting source figures.  This removes placeholder labels/dashes from
        # the layout background while keeping each real figure independently
        # editable as a PPTX picture object.
        replace_placeholders(
            base_image=base,
            spec=dict(spec),
            asset_dir=asset_dir,
            out_path=background,
            scale=scale,
            paste_assets=False,
        )

        prs = Presentation()
        prs.slide_width = Inches(slide_w_in)
        prs.slide_height = Inches(slide_h_in)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            str(background),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

        px_to_emu_x = prs.slide_width / width_px
        px_to_emu_y = prs.slide_height / height_px
        prepared = _prepare_figure_images(
            spec=spec,
            asset_dir=Path(asset_dir),
            out_dir=tmp,
            scale=scale,
        )
        for fig_id, figure_path, box in prepared:
            x0, y0, x1, y1 = box
            pic = slide.shapes.add_picture(
                str(figure_path),
                int(round(x0 * px_to_emu_x)),
                int(round(y0 * px_to_emu_y)),
                width=int(round((x1 - x0) * px_to_emu_x)),
                height=int(round((y1 - y0) * px_to_emu_y)),
            )
            pic.name = fig_id

        prs.save(out)
    return out


def _prepare_figure_images(
    *,
    spec: Mapping[str, Any],
    asset_dir: Path,
    out_dir: Path,
    scale: float,
) -> list[tuple[str, Path, tuple[int, int, int, int]]]:
    placements = dict(spec.get("placements") or {})
    placeholders = {str(p.get("id")): p for p in spec.get("placeholders", [])}
    rows: list[tuple[str, Path, tuple[int, int, int, int]]] = []
    for fig_id, raw_box in placements.items():
        ph = placeholders.get(str(fig_id), {})
        asset = ph.get("asset")
        if not asset:
            continue
        try:
            x0, y0, x1, y1 = [int(round(float(v) * scale)) for v in raw_box]
        except Exception:
            continue
        if x1 <= x0 or y1 <= y0:
            continue
        figure_path = out_dir / f"{_safe_name(str(fig_id))}.png"
        canvas = Image.new("RGB", (x1 - x0, y1 - y0), "white")
        paste_fit(canvas, asset_dir / str(asset), (0, 0, x1 - x0, y1 - y0), pad=0)
        canvas.save(figure_path, quality=95)
        rows.append((str(fig_id), figure_path, (x0, y0, x1, y1)))
    return rows


def _image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as im:
        width_px, height_px = im.size
    if width_px <= 0 or height_px <= 0:
        raise ValueError(f"cannot export PPTX from empty image: {path}")
    return width_px, height_px


def _slide_inches(width_px: int, height_px: int, *, long_edge_inches: float) -> tuple[float, float]:
    aspect = width_px / height_px
    if width_px >= height_px:
        slide_w_in = float(long_edge_inches)
        slide_h_in = slide_w_in / aspect
    else:
        slide_h_in = float(long_edge_inches)
        slide_w_in = slide_h_in * aspect
    return slide_w_in, slide_h_in


def _safe_name(value: str) -> str:
    keep = [ch if ch.isalnum() else "_" for ch in value]
    return "".join(keep).strip("_") or "figure"
