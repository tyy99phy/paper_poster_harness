from __future__ import annotations

import io
import re
from pathlib import Path
from zipfile import ZipFile
from typing import Iterable

from PIL import Image, ImageOps, ImageDraw

from .fonts import load_font


def extract_text(input_path: str | Path, out_txt: str | Path | None = None) -> str:
    p = Path(input_path)
    suffix = p.suffix.lower()
    if suffix == ".pdf":
        text = _extract_pdf_text(p)
    elif suffix == ".pptx":
        text = _extract_pptx_text(p)
    elif suffix in {".txt", ".md", ".tex"}:
        text = p.read_text(encoding="utf-8", errors="replace")
    else:
        raise ValueError(f"Unsupported text input type: {suffix}")
    text = normalize_text(text)
    if out_txt:
        Path(out_txt).parent.mkdir(parents=True, exist_ok=True)
        Path(out_txt).write_text(text, encoding="utf-8")
    return text


def _extract_pdf_text(path: Path) -> str:
    try:
        import fitz  # PyMuPDF
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("PyMuPDF/fitz is required for PDF text extraction") from exc
    doc = fitz.open(path)
    parts = []
    for i, page in enumerate(doc):
        parts.append(f"\n\n--- PAGE {i+1} ---\n")
        parts.append(page.get_text("text"))
    return "".join(parts)


def _extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("python-pptx is required for PPTX text extraction") from exc
    prs = Presentation(path)
    parts = []
    for si, slide in enumerate(prs.slides, start=1):
        parts.append(f"\n\n--- SLIDE {si} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = " ".join(shape.text.split())
                if t:
                    parts.append(t + "\n")
    return "".join(parts)


def normalize_text(text: str) -> str:
    lines = [re.sub(r"\s+", " ", line).strip() for line in text.splitlines()]
    lines = [line for line in lines if line]
    return "\n".join(lines) + "\n"


def extract_pptx_media(pptx_path: str | Path, out_dir: str | Path, contact_sheet: bool = True) -> list[Path]:
    pptx = Path(pptx_path)
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    paths: list[Path] = []
    with ZipFile(pptx) as zf:
        for name in zf.namelist():
            if not name.startswith("ppt/media/"):
                continue
            dest = out / Path(name).name
            dest.write_bytes(zf.read(name))
            paths.append(dest)
    if contact_sheet:
        make_contact_sheet(paths, out / "contact_sheet.jpg")
    return paths


def render_pdf_pages(pdf_path: str | Path, out_dir: str | Path, dpi: int = 200, max_pages: int | None = None) -> list[Path]:
    try:
        import fitz
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("PyMuPDF/fitz is required for PDF rendering") from exc
    pdf = Path(pdf_path)
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(pdf)
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    paths: list[Path] = []
    for i, page in enumerate(doc):
        if max_pages is not None and i >= max_pages:
            break
        pix = page.get_pixmap(matrix=mat, alpha=False)
        dest = out / f"page-{i+1:03d}.png"
        pix.save(dest)
        paths.append(dest)
    return paths


def extract_pdf_images(
    pdf_path: str | Path,
    out_dir: str | Path,
    *,
    min_width: int = 96,
    min_height: int = 96,
    max_images: int | None = None,
    contact_sheet: bool = True,
) -> list[Path]:
    """Extract raster images embedded in a PDF.

    This is intentionally conservative. It is a useful automatic source of
    replacement assets, but vector-only figures may not appear here; callers can
    still pass explicit asset directories such as ``figures/`` or
    ``figures_png/``.
    """
    try:
        import fitz
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("PyMuPDF/fitz is required for PDF image extraction") from exc
    pdf = Path(pdf_path)
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(pdf)
    paths: list[Path] = []
    seen_xrefs: set[int] = set()
    for page_index, page in enumerate(doc, start=1):
        for image_index, info in enumerate(page.get_images(full=True), start=1):
            if max_images is not None and len(paths) >= max_images:
                break
            xref = int(info[0])
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)
            try:
                extracted = doc.extract_image(xref)
                image_bytes = extracted.get("image")
                ext = str(extracted.get("ext") or "png").lower().replace("jpeg", "jpg")
                if not image_bytes:
                    continue
                # Validate dimensions before saving; this filters masks/icons.
                from PIL import Image

                with Image.open(io.BytesIO(image_bytes)) as im:
                    width, height = im.size
                    if width < min_width or height < min_height:
                        continue
                dest = out / f"pdf-page{page_index:03d}-img{image_index:03d}.{ext}"
                dest.write_bytes(image_bytes)
                paths.append(dest)
            except Exception:
                continue
        if max_images is not None and len(paths) >= max_images:
            break
    if contact_sheet:
        make_contact_sheet(paths, out / "contact_sheet.jpg")
    return paths


def make_contact_sheet(paths: Iterable[Path], out_path: str | Path, cell: tuple[int, int] = (320, 260), cols: int = 4) -> None:
    items = []
    for p in paths:
        try:
            im = Image.open(p).convert("RGB")
            items.append((p, im.copy()))
        except Exception:
            continue
    if not items:
        return
    rows = (len(items) + cols - 1) // cols
    w, h = cell
    sheet = Image.new("RGB", (cols * w, rows * h), "white")
    draw = ImageDraw.Draw(sheet)
    font = load_font(14)
    for idx, (p, im) in enumerate(items):
        x = (idx % cols) * w
        y = (idx // cols) * h
        thumb = ImageOps.contain(im, (w - 20, h - 55), Image.Resampling.LANCZOS)
        sheet.paste(thumb, (x + (w - thumb.width)//2, y + 8))
        draw.rectangle([x, y, x+w-1, y+h-1], outline=(220, 225, 235))
        draw.text((x+8, y+h-42), f"{p.name} {im.width}x{im.height}", fill=(0,0,0), font=font)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    sheet.save(out_path, quality=92)
