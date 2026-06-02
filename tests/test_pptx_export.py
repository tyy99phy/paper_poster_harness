from pathlib import Path

from PIL import Image
from pptx import Presentation

from poster_harness.pptx_export import export_image_as_pptx, export_layered_pptx


def _pictures(prs: Presentation):
    return [shape for shape in prs.slides[0].shapes if shape.shape_type == 13]


def test_export_image_as_pptx_preserves_portrait_aspect(tmp_path: Path):
    image = tmp_path / "poster.png"
    Image.new("RGB", (100, 150), "white").save(image)
    out = export_image_as_pptx(image, tmp_path / "poster.pptx")
    assert out.exists()
    prs = Presentation(out)
    assert len(prs.slides) == 1
    assert abs((prs.slide_width / prs.slide_height) - (100 / 150)) < 0.01
    pictures = _pictures(prs)
    assert len(pictures) == 1
    assert pictures[0].left == 0
    assert pictures[0].top == 0
    assert pictures[0].width == prs.slide_width
    assert pictures[0].height == prs.slide_height


def test_export_layered_pptx_keeps_each_source_figure_as_object(tmp_path: Path):
    base = tmp_path / "layout.png"
    Image.new("RGB", (200, 300), "white").save(base)
    asset_dir = tmp_path / "assets"
    asset_dir.mkdir()
    Image.new("RGB", (50, 50), "red").save(asset_dir / "fig1.png")
    Image.new("RGB", (80, 40), "blue").save(asset_dir / "fig2.png")
    spec = {
        "placeholders": [
            {"id": "FIG 01", "label": "Result", "aspect": "1:1 square", "asset": "fig1.png"},
            {"id": "FIG 02", "label": "Validation", "aspect": "2:1 wide", "asset": "fig2.png"},
        ],
        "placements": {
            "FIG 01": [20, 30, 80, 90],
            "FIG 02": [100, 120, 180, 160],
        },
        "_replacement_clear_boxes": {
            "FIG 01": [18, 28, 82, 92],
            "FIG 02": [98, 118, 182, 162],
        },
    }
    out = export_layered_pptx(
        base_image=base,
        spec=spec,
        asset_dir=asset_dir,
        out_path=tmp_path / "layered.pptx",
    )
    prs = Presentation(out)
    pictures = _pictures(prs)
    # One cleaned layout background plus two independently selectable figure objects.
    assert len(pictures) == 3
    assert [pic.name for pic in pictures[1:]] == ["FIG 01", "FIG 02"]
    assert pictures[0].left == 0
    assert pictures[0].top == 0
    assert pictures[0].width == prs.slide_width
    assert pictures[0].height == prs.slide_height
    assert pictures[1].left > 0
    assert pictures[1].top > 0
    assert pictures[2].left > pictures[1].left
